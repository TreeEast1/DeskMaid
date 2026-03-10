"""File content extraction for deep classification mode."""

from pathlib import Path

from deskmaid.scanner import FileInfo

MAX_CHARS = 500

READABLE_TEXT_SUFFIXES = {
    ".txt", ".md", ".csv", ".json", ".xml", ".yaml", ".yml",
    ".py", ".js", ".ts", ".java", ".c", ".cpp", ".h", ".go",
    ".rs", ".rb", ".php", ".sh", ".bat", ".ps1", ".sql",
    ".html", ".css", ".log", ".ini", ".cfg", ".conf", ".toml",
}

OFFICE_SUFFIXES = {".docx", ".xlsx", ".pptx"}


def read_text_file(path: Path) -> str | None:
    """Read first MAX_CHARS characters of a text file."""
    try:
        text = path.read_text(encoding="utf-8")
        return text[:MAX_CHARS]
    except UnicodeDecodeError:
        try:
            text = path.read_text(encoding="latin-1")
            return text[:MAX_CHARS]
        except Exception:
            return None
    except Exception:
        return None


def read_docx(path: Path) -> str | None:
    """Extract text from a .docx file."""
    try:
        from docx import Document
        doc = Document(str(path))
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return text[:MAX_CHARS] if text else None
    except Exception:
        return None


def read_xlsx(path: Path) -> str | None:
    """Extract text from a .xlsx file (first sheet headers + first rows)."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        ws = wb.active
        lines = []
        for row in ws.iter_rows(max_row=10, values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append(" | ".join(cells))
        wb.close()
        text = "\n".join(lines)
        return text[:MAX_CHARS] if text else None
    except Exception:
        return None


def read_pptx(path: Path) -> str | None:
    """Extract text from a .pptx file."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text.strip())
        text = "\n".join(texts)
        return text[:MAX_CHARS] if text else None
    except Exception:
        return None


def read_pdf(path: Path) -> str | None:
    """Extract text from first 2 pages of a PDF file."""
    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(str(path)) as pdf:
            for page in pdf.pages[:2]:
                page_text = page.extract_text()
                if page_text:
                    texts.append(page_text)
        text = "\n".join(texts)
        return text[:MAX_CHARS] if text else None
    except Exception:
        return None


def extract_content(file_info: FileInfo) -> str | None:
    """Extract content preview from a file based on its suffix."""
    if file_info.is_dir:
        return None

    suffix = file_info.suffix.lower()
    path = file_info.path

    if suffix in READABLE_TEXT_SUFFIXES:
        return read_text_file(path)
    if suffix == ".docx":
        return read_docx(path)
    if suffix == ".xlsx":
        return read_xlsx(path)
    if suffix == ".pptx":
        return read_pptx(path)
    if suffix == ".pdf":
        return read_pdf(path)

    return None


def enrich_items_with_content(items: list[FileInfo]) -> list[dict]:
    """Add content_preview field to each item's dict representation."""
    enriched = []
    for item in items:
        d = item.to_dict()
        content = extract_content(item)
        if content:
            d["content_preview"] = content
        enriched.append(d)
    return enriched
